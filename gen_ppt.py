#!/usr/bin/env python3
"""Génère la présentation PPT KopéAgri Caraïbes pour partenariats"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Couleurs
GREEN_DARK = RGBColor(0x1B, 0x5E, 0x20)
GREEN_MED = RGBColor(0x2E, 0x7D, 0x32)
GREEN_LIGHT = RGBColor(0x4C, 0xAF, 0x50)
GREEN_BG = RGBColor(0xE8, 0xF5, 0xE9)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x21, 0x21, 0x21)
GRAY = RGBColor(0x75, 0x75, 0x75)
GOLD = RGBColor(0xFF, 0xA0, 0x00)
WHATSAPP = RGBColor(0x25, 0xD3, 0x66)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def add_bg(slide, color=GREEN_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text_box(slide, left, top, width, height, text, font_size=18, bold=False, color=WHITE, alignment=PP_ALIGN.LEFT, font_name='Calibri'):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_shape(slide, left, top, width, height, fill_color, text='', font_size=14, font_color=WHITE):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = font_color
        p.font.bold = True
        p.font.name = 'Calibri'
    return shape

# ===== SLIDE 1: COVER =====
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
add_bg(slide, GREEN_DARK)

# Dégradé simulé avec rectangle
add_shape(slide, 0, 0, 13.333, 7.5, GREEN_DARK)

# Titre principal
add_text_box(slide, 1, 1.5, 11, 1.5, '🌴 KopéAgri Caraïbes', font_size=52, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

# Sous-titre
add_text_box(slide, 1.5, 3.2, 10, 1, 'La plateforme coopérative agricole digitale\nde Martinique et des Caraïbes', font_size=28, color=GREEN_LIGHT, alignment=PP_ALIGN.CENTER)

# Bas
add_text_box(slide, 2, 5.5, 9, 0.6, 'Développons ensemble le réseau agricole caraïbéen', font_size=20, color=GOLD, alignment=PP_ALIGN.CENTER)
add_text_box(slide, 2, 6.3, 9, 0.5, 'contact@delikreol.mq  •  cvlad97.github.io/kopeagri-caraibes', font_size=16, color=RGBColor(0xA5, 0xD6, 0xA7), alignment=PP_ALIGN.CENTER)

# ===== SLIDE 2: PROBLÈME =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_text_box(slide, 0.5, 0.3, 12, 0.8, 'Le constat en Martinique', font_size=36, bold=True, color=GREEN_DARK, alignment=PP_ALIGN.CENTER)

problems = [
    ('📉', 'Producteurs isolés', 'Difficulté à écouler la production, manque de visibilité et d\'accès aux marchés'),
    ('🚛', 'Transport deficient', 'Pas de coordination logistique entre les zones rurales et les points de vente'),
    ('📱', 'Fracture numérique', 'Outils digitaux inadaptés pour une population agricole non-tech'),
    ('🔄', 'Pas de mise en relation', 'Aucun système centralisé pour connecter producteurs, transporteurs et acheteurs'),
]

for i, (emoji, title, desc) in enumerate(problems):
    y = 1.5 + i * 1.4
    add_shape(slide, 0.8, y, 11.5, 1.2, GREEN_BG)
    add_text_box(slide, 1, y + 0.1, 0.8, 0.8, emoji, font_size=36, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, 2, y + 0.1, 3, 0.5, title, font_size=22, bold=True, color=GREEN_DARK)
    add_text_box(slide, 2, y + 0.6, 9.5, 0.5, desc, font_size=16, color=GRAY)

# ===== SLIDE 3: SOLUTION =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, GREEN_DARK)

add_text_box(slide, 0.5, 0.3, 12, 0.8, 'KopéAgri : La solution', font_size=36, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

solutions = [
    ('👨‍🌾', 'Fiches Producteurs', 'Répertoire digital avec photos, cultures, certifications et contact direct'),
    ('🚛', 'Mise en relation Transport', 'Matching géographique automatique entre producteurs et transporteurs'),
    ('📤', 'Appels d\'Offre Rapides', 'Création RFQ en 60s, matching partenaires, suivi et confirmation via WhatsApp'),
    ('💬', 'WhatsApp Natif', 'Communication intégrée — chaque action génère un message WhatsApp pré-rempli'),
    ('📊', 'Dashboard par Rôle', 'Interface adaptée : producteur, transporteur, acheteur, coopérative, institution'),
]

for i, (emoji, title, desc) in enumerate(solutions):
    y = 1.3 + i * 1.15
    add_text_box(slide, 1, y, 0.7, 0.7, emoji, font_size=30, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, 2, y + 0.05, 3.5, 0.5, title, font_size=20, bold=True, color=GOLD)
    add_text_box(slide, 5.5, y + 0.05, 7, 0.5, desc, font_size=16, color=RGBColor(0xC8, 0xE6, 0xC9))

# ===== SLIDE 4: COMMENT ÇA MARCHE =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_text_box(slide, 0.5, 0.3, 12, 0.8, 'Comment ça marche ?', font_size=36, bold=True, color=GREEN_DARK, alignment=PP_ALIGN.CENTER)

steps = [
    ('1', 'Inscription', '3 étapes simples\nPhoto + Nom + Téléphone', GREEN_LIGHT),
    ('2', 'Appel d\'Offre', 'Décrivez votre besoin\nProduit, quantité, date', GOLD),
    ('3', 'Matching Auto', '5 partenaires trouvés\npar proximité géographique', RGBColor(0x42, 0xA5, 0xF5)),
    ('4', 'WhatsApp', 'Contact direct\nMessage pré-rempli', WHATSAPP),
    ('5', 'Suivi', 'Pipeline en temps réel\nConfirmé → En cours → Livré', RGBColor(0xAB, 0x47, 0xBC)),
]

for i, (num, title, desc, color) in enumerate(steps):
    x = 0.5 + i * 2.5
    add_shape(slide, x, 1.8, 2.2, 4.5, color)
    add_text_box(slide, x, 2.0, 2.2, 0.8, num, font_size=48, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + 0.1, 2.9, 2, 0.6, title, font_size=20, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + 0.1, 3.6, 2, 2, desc, font_size=14, color=WHITE, alignment=PP_ALIGN.CENTER)

    # Flèche entre les étapes
    if i < len(steps) - 1:
        add_text_box(slide, x + 2.2, 3.5, 0.3, 0.5, '→', font_size=28, bold=True, color=GREEN_DARK, alignment=PP_ALIGN.CENTER)

# ===== SLIDE 5: CHIFFRES =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, GREEN_DARK)

add_text_box(slide, 0.5, 0.3, 12, 0.8, 'La Martinique agricole en chiffres', font_size=36, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

chiffres = [
    ('7 000+', 'Exploitations\nagricoles'),
    ('34', 'Communes\ncouvertes'),
    ('27', 'Cultures\nlocales'),
    ('100%', 'Conçu pour\nles Caraïbes'),
]

for i, (num, label) in enumerate(chiffres):
    x = 1 + i * 3
    add_shape(slide, x, 2, 2.5, 3.5, RGBColor(0x2E, 0x7D, 0x32))
    add_text_box(slide, x, 2.3, 2.5, 1.2, num, font_size=52, bold=True, color=GOLD, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x, 3.8, 2.5, 1.2, label, font_size=18, color=WHITE, alignment=PP_ALIGN.CENTER)

add_text_box(slide, 1, 6, 11, 0.8, 'Banane, canne, ananas, mangue, igname, dasheen, piment, cacao, café, vanille...', font_size=16, color=RGBColor(0xA5, 0xD6, 0xA7), alignment=PP_ALIGN.CENTER)

# ===== SLIDE 6: APPEL D'OFFRE EN DÉTAIL =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_text_box(slide, 0.5, 0.3, 12, 0.8, 'Le système d\'Appel d\'Offre', font_size=36, bold=True, color=GREEN_DARK, alignment=PP_ALIGN.CENTER)

# Pipeline
pipeline = [
    ('📝 Brouillon', RGBColor(0x9E, 0x9E, 0x9E)),
    ('📤 Envoyée', RGBColor(0x21, 0x96, 0xF3)),
    ('✅ Confirmée', RGBColor(0x4C, 0xAF, 0x50)),
    ('🔄 En cours', RGBColor(0xFF, 0x98, 0x00)),
    ('📦 Livrée', RGBColor(0x1B, 0x5E, 0x20)),
]

for i, (label, color) in enumerate(pipeline):
    x = 0.5 + i * 2.5
    add_shape(slide, x, 1.5, 2.2, 0.8, color, label, font_size=16, font_color=WHITE)
    if i < len(pipeline) - 1:
        add_text_box(slide, x + 2.2, 1.55, 0.3, 0.5, '→', font_size=24, bold=True, color=GREEN_DARK, alignment=PP_ALIGN.CENTER)

# Types d'appels
types_data = [
    ('🚛 Transport', 'Producteur → Transporteur\nFrigorifique, volume, trajet'),
    ('🥭 Achat / Vente', 'Producteur ↔ Acheteur B2B\nPrix, quantité, qualité'),
    ('🏠 Stockage', 'Producteur → Dépositaire\nCapacité, durée, conditions'),
    ('🌍 Export', 'Producteur → Exportateur\nCertifications, volumes, Incoterms'),
]

for i, (title, desc) in enumerate(types_data):
    x = 0.5 + (i % 2) * 6.2
    y = 3 + (i // 2) * 2
    add_shape(slide, x, y, 5.8, 1.6, GREEN_BG)
    add_text_box(slide, x + 0.2, y + 0.1, 5.4, 0.5, title, font_size=20, bold=True, color=GREEN_DARK)
    add_text_box(slide, x + 0.2, y + 0.7, 5.4, 0.7, desc, font_size=14, color=GRAY)

# ===== SLIDE 7: AVANTAGES PARTENAIRES =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, GREEN_DARK)

add_text_box(slide, 0.5, 0.3, 12, 0.8, 'Pourquoi devenir partenaire ?', font_size=36, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

avantages = [
    ('🏪', 'Acheteurs / Distributeurs', 'Accès direct aux producteurs locaux\nQualité traçable, prix transparents\nAppels d\'offre automatiques sur vos besoins'),
    ('🚛', 'Transporteurs', 'Flux régulier de missions\nOptimisation des trajets\nConfirmation rapide via WhatsApp'),
    ('🏛️', 'Institutions / Collectivités', 'Visibilité sur la production locale\nDonnées pour les politiques agricoles\nRelais territorial dans les 34 communes'),
    ('🤝', 'Coopératives / ONG', 'Outil digital pour vos adhérents\nCoordination logistique partagée\nReporting et suivi automatisés'),
]

for i, (emoji, title, desc) in enumerate(avantages):
    x = 0.5 + (i % 2) * 6.2
    y = 1.3 + (i // 2) * 2.8
    add_shape(slide, x, y, 5.8, 2.4, RGBColor(0x2E, 0x7D, 0x32))
    add_text_box(slide, x + 0.3, y + 0.2, 0.8, 0.6, emoji, font_size=32, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + 1.2, y + 0.2, 4, 0.5, title, font_size=20, bold=True, color=GOLD)
    add_text_box(slide, x + 1.2, y + 0.8, 4.2, 1.4, desc, font_size=14, color=RGBColor(0xC8, 0xE6, 0xC9))

# ===== SLIDE 8: DEMANDE =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_text_box(slide, 0.5, 0.3, 12, 0.8, 'Ce que nous cherchons', font_size=36, bold=True, color=GREEN_DARK, alignment=PP_ALIGN.CENTER)

demandes = [
    ('🔗', 'Réseau', 'Nous connecter avec des producteurs, transporteurs, acheteurs et institutions dans toute la Caraïbe'),
    ('📡', 'Visibilité', 'Relayer KopéAgri auprès de votre réseau : associations, chambres d\'agriculture, médias locaux'),
    ('💼', 'Partenariats', 'Co-construire des appels d\'offre pilotes dans votre secteur géographique ou thématique'),
    ('🧪', 'Testeurs', 'Agriculteurs et opérateurs prêts à tester la plateforme et nous remonter les améliorations'),
]

for i, (emoji, title, desc) in enumerate(demandes):
    y = 1.5 + i * 1.4
    add_shape(slide, 0.8, y, 11.5, 1.2, GREEN_BG)
    add_text_box(slide, 1, y + 0.1, 0.8, 0.8, emoji, font_size=36, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, 2, y + 0.1, 3, 0.5, title, font_size=22, bold=True, color=GREEN_DARK)
    add_text_box(slide, 2, y + 0.6, 9.5, 0.5, desc, font_size=16, color=GRAY)

# ===== SLIDE 9: CTA =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, GREEN_DARK)

add_text_box(slide, 1, 1, 11, 1.2, 'Rejoignez le mouvement 🌴', font_size=48, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

add_text_box(slide, 2, 2.8, 9, 1, 'KopéAgri Caraïbes a besoin de vous pour construire\nle réseau agricole digital de demain', font_size=24, color=RGBColor(0xA5, 0xD6, 0xA7), alignment=PP_ALIGN.CENTER)

# Box lien
add_shape(slide, 3, 4.2, 7, 1.2, WHATSAPP, '💬 Écrivez-nous sur WhatsApp', font_size=22, font_color=WHITE)
add_text_box(slide, 3, 5.5, 7, 0.5, 'wa.me/596696000000', font_size=18, color=RGBColor(0xA5, 0xD6, 0xA7), alignment=PP_ALIGN.CENTER)

add_shape(slide, 3, 5.8, 7, 0.6, RGBColor(0x2E, 0x7D, 0x32), '🌐 cvlad97.github.io/kopeagri-caraibes', font_size=16, font_color=WHITE)

add_text_box(slide, 2, 6.7, 9, 0.5, 'contact@delikreol.mq  •  Ensemble, faisons KopéAgri !', font_size=16, color=RGBColor(0x81, 0xC7, 0x84), alignment=PP_ALIGN.CENTER)

# Save
output = '/workspace/kopeagri-caraibes/KopeAgri_Caraibes_Partenariats.pptx'
prs.save(output)
print(f'PPT saved: {output}')
